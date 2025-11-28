import logging
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from sqlalchemy.orm import Session
import os
import shutil
import asyncio
import tempfile
from datetime import datetime

from models.database import get_db, Document, User
from api.auth.jwt_handler import get_current_active_user
from .schemas import DocumentResponse, DocumentList
from api.features.rag_manager import get_rag_instance, reload_rag_instance
from PathRAG.utils import compute_mdhash_id

# Additional libraries for file processing
import PyPDF2
import docx2txt
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load as load_odf
from odf import teletype
from ebooklib import epub
from bs4 import BeautifulSoup
logger = logging.getLogger("PathRAG")
# Create uploads directory if it doesn't exist
UPLOAD_DIR = "./uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

router = APIRouter(
    prefix="/documents",
    tags=["Documents"],
    dependencies=[Depends(get_current_active_user)]
)

# No WebSocket or upload status tracking needed

# Setup a working directory for PathRAG.
WORKING_DIR = os.path.join(os.getcwd(), 'data')
if not os.path.exists(WORKING_DIR):
    os.mkdir(WORKING_DIR)

# Get the RAG instance from the manager
rag = get_rag_instance()

def extract_text_from_file(file: UploadFile) -> str:
    """
    Extract text from an uploaded file.
    Supports many file types including:
    .txt, .md, .pdf, .docx, .pptx, .xlsx, .rtf, .odt, .tex, .epub,
    .html, .htm, .csv, .json, .xml, .yaml, .yml, .log, .conf, .ini,
    .properties, .sql, .bat, .sh, .c, .cpp, .py, .java, .js, .ts,
    .swift, .go, .rb, .php, .css, .scss, .less.
    """
    filename = file.filename
    extension = os.path.splitext(filename)[1].lower()
    file.file.seek(0)

    # Define plain text file extensions.
    plain_text_ext = [
        ".txt", ".md", ".tex", ".csv", ".json", ".xml", ".yaml", ".yml",
        ".log", ".conf", ".ini", ".properties", ".sql", ".bat", ".sh",
        ".c", ".cpp", ".py", ".java", ".js", ".ts", ".swift", ".go",
        ".rb", ".php", ".css", ".scss", ".less"
    ]

    if extension in plain_text_ext:
        try:
            return file.file.read().decode('utf-8', errors='ignore')
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error reading {extension} file: {str(e)}")
    elif extension == ".pdf":
        try:
            file.file.seek(0)
            pdf_reader = PyPDF2.PdfReader(file.file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing PDF file: {str(e)}")
    elif extension == ".docx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            text = docx2txt.process(tmp_path)
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing DOCX file: {str(e)}")
    elif extension == ".pptx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            prs = Presentation(tmp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing PPTX file: {str(e)}")
    elif extension == ".xlsx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    text += row_text + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing XLSX file: {str(e)}")
    elif extension == ".rtf":
        try:
            file.file.seek(0)
            content = file.file.read().decode('utf-8', errors='ignore')
            return rtf_to_text(content)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing RTF file: {str(e)}")
    elif extension == ".odt":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".odt") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            doc = load_odf(tmp_path)
            text_content = teletype.extractText(doc)
            os.remove(tmp_path)
            return text_content
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing ODT file: {str(e)}")
    elif extension == ".epub":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".epub") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            book = epub.read_epub(tmp_path)
            text = ""
            for item in book.get_items():
                if item.get_type() == epub.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text += soup.get_text() + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing EPUB file: {str(e)}")
    elif extension in [".html", ".htm"]:
        try:
            file.file.seek(0)
            content = file.file.read().decode('utf-8', errors='ignore')
            soup = BeautifulSoup(content, "html.parser")
            return soup.get_text()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing HTML file: {str(e)}")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type.")

@router.post("/upload", response_model=DocumentResponse, summary="Upload a file", description="Upload a file to insert its content into the RAG system.")
async def upload_file(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user)
):
    try:
        # Validate file type
        filename = file.filename
        file_extension = os.path.splitext(filename)[1].lower()

        # Map file extensions to content types
        content_type_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.md': 'text/markdown',
            '.txt': 'text/plain',
            '.html': 'text/html',
            '.htm': 'text/html',
            '.rtf': 'application/rtf',
            '.odt': 'application/vnd.oasis.opendocument.text',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.epub': 'application/epub+zip'
        }

        content_type = content_type_map.get(file_extension, file.content_type)

        # Create file path
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        file_path = os.path.join(UPLOAD_DIR, f"{current_user.id}_{timestamp}_{filename}")

        # Save file
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Get file size
        file_size = os.path.getsize(file_path)

        # Create document record
        db_document = Document(
            user_id=current_user.id,
            filename=filename,
            content_type=content_type,
            file_path=file_path,
            file_size=file_size,
            status="processing"
        )
        db.add(db_document)
        db.commit()
        db.refresh(db_document)

        # Extract text and process with PathRAG in background
        async def process_document_task():
            # Declare global rag at the beginning of the function
            global rag

            try:
                # Extract text from file
                with open(file_path, "rb") as f:
                    # Create a temporary UploadFile-like object
                    class TempUploadFile:
                        def __init__(self, file_path):
                            self.filename = os.path.basename(file_path)
                            self.file = open(file_path, "rb")

                    temp_file = TempUploadFile(file_path)
                    content = extract_text_from_file(temp_file)
                    temp_file.file.close()

                # Process with PathRAG
                await rag.ainsert(content)

                # Reload the PathRAG instance to make the new document available
                reload_rag_instance()

                # Update the local reference
                rag = get_rag_instance()

                # Update document status
                document = db.query(Document).filter(Document.id == db_document.id).first()
                if document:
                    document.status = "completed"
                    document.processed_at = datetime.now()
                    db.commit()

                logger.info(f"Document processed and PathRAG reloaded successfully.")
            except Exception as e:
                # Update document status on error
                document = db.query(Document).filter(Document.id == db_document.id).first()
                if document:
                    document.status = "failed"
                    document.error_message = str(e)
                    db.commit()

        # Start processing task
        task = asyncio.create_task(process_document_task())

        # Add task to a global set to prevent garbage collection
        if not hasattr(upload_file, 'tasks'):
            upload_file.tasks = set()
        upload_file.tasks.add(task)
        task.add_done_callback(lambda t: upload_file.tasks.remove(t))
        logger.info(f"Document upload completed.")
        return db_document
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/", response_model=DocumentList)
async def get_documents(db: Session = Depends(get_db), current_user: User = Depends(get_current_active_user)):
    documents = db.query(Document).filter(Document.user_id == current_user.id).all()
    return {"documents": documents}

@router.get("/{document_id}", response_model=DocumentResponse)
async def get_document(document_id: int, db: Session = Depends(get_db), current_user: User = Depends(get_current_active_user)):
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if document is None:
        raise HTTPException(status_code=404, detail="Document not found")
    return document

@router.get("/{document_id}/status", response_model=DocumentResponse)
async def get_document_status(document_id: int, db: Session = Depends(get_db), current_user: User = Depends(get_current_active_user)):
    # Check if document exists and belongs to user
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if document is None:
        raise HTTPException(status_code=404, detail="Document not found")

    # Return the document with its current status
    return document

@router.delete("/{document_id}", summary="Delete a document", description="Delete a document and its associated graph elements.")
async def delete_document(
    document_id: int, 
    db: Session = Depends(get_db), 
    current_user: User = Depends(get_current_active_user)
):
    global rag
    
    # Check if document exists and belongs to user
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if document is None:
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Delete from RAG
    try:
        # Read content to compute ID
        if os.path.exists(document.file_path):
             # Re-open file
             with open(document.file_path, "rb") as f:
                  # Create temp file object similar to upload_file logic
                   class TempUploadFile:
                        def __init__(self, file_path):
                            self.filename = os.path.basename(file_path)
                            self.file = open(file_path, "rb")
                   
                   temp_file = TempUploadFile(document.file_path)
                   content = extract_text_from_file(temp_file)
                   temp_file.file.close()
             
             doc_hash_id = compute_mdhash_id(content.strip(), prefix="doc-")
             await rag.adelete_document(doc_hash_id)
             
             # Also remove file from disk
             try:
                 os.remove(document.file_path)
             except Exception as e:
                 logger.warning(f"Failed to delete file from disk: {e}")
        else:
            logger.warning(f"File {document.file_path} not found on disk.")
            # Still try to delete from RAG if we could compute hash? 
            # Without content, we can't compute hash. So we can't delete from RAG if file is missing.
            # But maybe we should delete from DB anyway.
            pass

        # Remove from DB
        db.delete(document)
        db.commit()
        
        return {"message": "Document deleted successfully"}
        
    except Exception as e:
        logger.error(f"Error deleting document: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/reload", response_model=dict)
async def reload_documents(current_user: User = Depends(get_current_active_user)):
    """
    Reload the PathRAG instance to recognize newly uploaded documents.
    This is useful after uploading documents to make them available for querying
    without restarting the server.
    """
    # Declare global rag at the beginning of the function
    global rag

    try:
        # Reload the PathRAG instance
        reload_rag_instance()

        # Update the local reference
        rag = get_rag_instance()

        return {
            "success": True,
            "message": "PathRAG instance reloaded successfully. New documents are now available for querying."
        }
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to reload PathRAG instance: {str(e)}"
        )
